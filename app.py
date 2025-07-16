import streamlit as st
import os
import json
import requests
from dotenv import load_dotenv
import io
import base64 # For decoding base64 image data
import asyncio # For running async functions

# Libraries for file parsing
import docx
from PyPDF2 import PdfReader # Using PyPDF2 for PDF reading
from youtube_transcript_api._api import YouTubeTranscriptApi
from youtube_transcript_api._errors import NoTranscriptFound, TranscriptsDisabled # For YouTube

# Langchain imports
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain_core.prompts import ChatPromptTemplate
from langchain_core.output_parsers import JsonOutputParser
from langchain_core.pydantic_v1 import BaseModel, Field
import google.generativeai as genai # For multimodal input

# python-pptx for PPTX generation
from pptx import Presentation as PPTXPresentation
from pptx.util import Inches
from pptx.enum.text import PP_ALIGN
#from pptx.enum.shapes import MSO_PLACEHOLDER_TYPE
#from pptx.enum.shapes import MSO_SHAPE_TYPE # Import MSO_SHAPE_TYPE for checking placeholder type
# Define MSO_PLACEHOLDER_TYPE manually if not available in pptx.enum.shapes
#class MSO_PLACEHOLDER_TYPE:
#    TITLE = 1
#    BODY = 2

# Ensure MSO_PLACEHOLDER_TYPE is available. If direct import fails, use manual definition.
try:
    from pptx.enum.shapes import MSO_PLACEHOLDER_TYPE
except ImportError:
    # Fallback for older python-pptx versions or specific environments
    class MSO_PLACEHOLDER_TYPE:
        TITLE = 1
        BODY = 2
        

# Load environment variables (for API key)
load_dotenv()

# --- Configuration ---
GOOGLE_API_KEY = os.getenv("GOOGLE_API_KEY")
if not GOOGLE_API_KEY:
    st.error("Google API Key not found. Please set the GOOGLE_API_KEY environment variable in a .env file.")
    st.stop()

UNSPLASH_ACCESS_KEY = os.getenv("UNSPLASH_ACCESS_KEY")
if not UNSPLASH_ACCESS_KEY:
    st.error("Unsplash Access Key not found. Please set the UNSPLASH_ACCESS_KEY environment variable in a .env file.")
    st.stop()

# Configure Google Generative AI globally for both direct calls and Langchain
genai.configure(api_key=GOOGLE_API_KEY)

# --- Pydantic Model for Structured Output ---
class Slide(BaseModel):
    title: str = Field(description="The concise title of the presentation slide.")
    bullet_points: list[str] = Field(description="A list of 3-5 key bullet points for the slide content.")

class Presentation(BaseModel):
    slides: list[Slide] = Field(description="A list of presentation slides, each with a title and bullet points.")

# --- File Reading Functions ---
def read_txt(file):
    """Reads content from a .txt file."""
    return file.read().decode("utf-8")

def read_docx(file):
    """Reads content from a .docx file."""
    doc = docx.Document(file)
    full_text = []
    for para in doc.paragraphs:
        full_text.append(para.text)
    return "\n".join(full_text)

def read_pdf(file):
    """Reads content from a .pdf file."""
    reader = PdfReader(file)
    full_text = []
    for page in reader.pages:
        full_text.append(page.extract_text())
    return "\n".join(full_text)

def get_youtube_transcript(youtube_url):
    """
    Fetches the transcript from a YouTube video URL.
    Returns the concatenated transcript text or None if not available/error.
    """
    try:
        video_id = youtube_url.split("v=")[1].split("&")[0]
        transcript_list = YouTubeTranscriptApi.list_transcripts(video_id)
        
        # Try to get English transcript first, fallback to any available
        transcript = None
        try:
            transcript = transcript_list.find_transcript(['en'])
        except NoTranscriptFound:
            transcript = None
            # Fallback: Try to get the first available transcript in any language
            for t in transcript_list:
                transcript = t
                break
        
        if transcript:
            full_transcript_data = transcript.fetch()
            # FIX: Access the 'text' attribute directly from the FetchedTranscriptSnippet object
            return " ".join([entry.text for entry in full_transcript_data])
        else:
            st.warning("No transcript found for the provided YouTube video in English or any other language.")
            return None
    except NoTranscriptFound:
        st.warning("No transcript found for this YouTube video. It might not have captions or they are disabled.")
        return None
    except TranscriptsDisabled:
        st.warning("Transcripts are disabled for this YouTube video.")
        return None
    except Exception as e:
        st.error(f"Error fetching YouTube transcript: {e}. Please check the URL and video availability.")
        return None

def get_image_description(image_bytes):
    """Generates a text description of an image using a multimodal model."""
    try:
        multimodal_model = genai.GenerativeModel('gemini-1.5-flash')
        image_part = {
            "mime_type": "image/jpeg", # Assuming common image types, adjust if needed
            "data": image_bytes
        }
        response = multimodal_model.generate_content(["Describe this image in detail:", image_part])
        return response.text
    except Exception as e:
        st.warning(f"Could not describe image: {e}. Proceeding without image description.")
        return ""

# --- Functions for AI Interaction ---

async def generate_slide_content(topic, num_slides, language, tone, audience, scene):
    """
    Generates presentation slide content using Google Gemini via Langchain.
    Returns a list of dictionaries, each representing a slide.
    """
    llm = ChatGoogleGenerativeAI(
        model="gemini-2.0-flash",
        client_options=None,
        transport=None,
        additional_headers=None,
        client=None,
        async_client=None
    )
    parser = JsonOutputParser(pydantic_object=Presentation)

    prompt = ChatPromptTemplate.from_messages(
        [
            ("system", "You are an expert presentation content generator. Create a presentation outline based on the user's input, formatted as a JSON array of slides."),
            ("human", """
            Generate a {num_slides}-slide presentation outline on the topic: "{topic}".
            Each slide should have a concise title and 3-5 bullet points.
            The presentation should be in {language}, with a {tone} tone, for an {audience} audience, and styled for a {scene} scene.

            {format_instructions}
            """)
        ]
    )

    chain = prompt | llm | parser

    try:
        response = await chain.ainvoke({
            "topic": topic,
            "num_slides": num_slides,
            "language": language,
            "tone": tone,
            "audience": audience,
            "scene": scene,
            "format_instructions": parser.get_format_instructions()
        })
        return response.get('slides', [])
    except Exception as e:
        st.error(f"Error generating slide content: {e}")
        return []

def generate_image_for_slide(slide_title):
    """
    Fetches an image from Unsplash based on the slide title.
    Returns the base64 encoded image data URL or a placeholder.
    """
    keywords = slide_title.split(':')[-1].strip() if ':' in slide_title else slide_title
    search_query = f"{keywords} professional abstract concept" # Refine search for better results

    # Unsplash API endpoint for searching photos
    unsplash_api_url = f"https://api.unsplash.com/search/photos?query={encodeURIComponent(search_query)}&per_page=1&client_id={UNSPLASH_ACCESS_KEY}"

    try:
        # First, search Unsplash for an image URL
        response = requests.get(unsplash_api_url)
        response.raise_for_status() # Raise an exception for HTTP errors (4xx or 5xx)
        data = response.json()

        if data.get("results") and len(data["results"]) > 0:
            image_url = data["results"][0]["urls"]["regular"] # Get the regular size image URL

            # Then, fetch the image content as bytes
            image_response = requests.get(image_url)
            image_response.raise_for_status()
            image_bytes = image_response.content

            # Base64 encode the image bytes
            base64_image = base64.b64encode(image_bytes).decode('utf-8')
            return f"data:image/png;base64,{base64_image}"
        else:
            st.warning(f"Unsplash API did not return an image for slide: '{slide_title}'. Using placeholder.")
            return f"https://placehold.co/600x400/E0E7FF/4338CA?text=Image+for+Slide"
    except requests.exceptions.RequestException as e:
        st.error(f"Error fetching image from Unsplash for slide '{slide_title}': {e}. Using placeholder.")
        return f"https://placehold.co/600x400/E0E7FF/4338CA?text=Image+Error"
    except Exception as e:
        st.error(f"An unexpected error occurred while processing Unsplash image for '{slide_title}': {e}. Using placeholder.")
        return f"https://placehold.co/600x400/E0E7FF/4338CA?text=Image+Error"

from urllib.parse import quote

def encodeURIComponent(s):
    """
    Encodes a string for use in a URL component, similar to JavaScript's encodeURIComponent.
    """
    return quote(s, safe='')

def create_pptx_file(slides_data, theme_style="Default"): # Added theme_style parameter
    """
    Creates a PowerPoint presentation (.pptx) from generated slide data.
    Returns a BytesIO object containing the PPTX file.
    """
    # --- Theme/Template Loading Logic ---
    template_path = None

    # FIX: Make base_template_dir absolute using the script's directory
    base_template_dir = os.path.join(os.path.dirname(__file__), "templates")

    if theme_style == "Classic":
        template_filename = "classic_template.pptx"
    elif theme_style == "Minimalist":
        template_filename = "minimalist_template.pptx"
    elif theme_style == "Organic":
        template_filename = "Organic.pptx"
    elif theme_style == "Tech Design":
        template_filename = "Tech design.pptx"
    elif theme_style == "Basic":
        template_filename = "Basic.pptx" # Assuming basic.pptx is the primary basic template
    else: # Default theme
        template_filename = None # No specific template, will use blank presentation

    if template_filename:
        template_path = os.path.join(base_template_dir, template_filename)
        if os.path.exists(template_path):
            prs = PPTXPresentation(template_path)
            st.success(f"Loaded '{theme_style}' theme template from {template_path}")
        else:
            prs = PPTXPresentation() # Fallback to blank presentation
            st.warning(f"Could not find template for '{theme_style}' theme at '{template_path}'. Creating a default blank presentation.")
    else:
        prs = PPTXPresentation() # Creates a new, blank presentation for "Default" theme

    # ... (rest of your create_pptx_file function remains the same)


    # Using a more robust way to find the title and body placeholders
    # These are standard slide layouts in python-pptx, which should generally work
    # if no custom template is loaded or if the custom template has these layouts.
    title_layout = prs.slide_layouts[0] # Typically Title Slide
    content_layout = prs.slide_layouts[1] # Typically Title and Content

    for i, slide in enumerate(slides_data):
        slide_obj = prs.slides.add_slide(content_layout)

        # Find the title placeholder
        title_shape = None
        for shape in slide_obj.shapes:
            if shape.has_text_frame and shape.is_placeholder and shape.placeholder_format.type == MSO_PLACEHOLDER_TYPE.TITLE:
                title_shape = shape
                break

        if title_shape:
            title_shape.text = slide['title']
        else:
            # Fallback if no title placeholder found, add a text box
            left = Inches(1)
            top = Inches(0.5)
            width = Inches(8)
            height = Inches(1)
            txBox = slide_obj.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            p = tf.add_paragraph()
            p.text = slide['title']
            p.alignment = PP_ALIGN.CENTER
            st.warning(f"No title placeholder found for slide '{slide['title']}'. Added title as a text box.")


        # Find the body (content) placeholder
        body_shape = None
        for shape in slide_obj.shapes:
            if shape.has_text_frame and shape.is_placeholder and shape.placeholder_format.type == MSO_PLACEHOLDER_TYPE.BODY:
                body_shape = shape
                break

        if body_shape and hasattr(body_shape, "has_text_frame") and body_shape.has_text_frame:
            tf = body_shape.text_frame
            tf.clear()

            for point in slide['bullet_points']:
                p = tf.add_paragraph()
                p.text = point
                p.level = 0
        else:
            # Fallback if no body placeholder found, add a text box for bullet points
            # Adjusted position to reduce overlap if no template placeholder is found
            left = Inches(1)
            top = Inches(2)
            width = Inches(8)
            height = Inches(4) # Keep height for text
            txBox = slide_obj.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.clear()
            for point in slide['bullet_points']:
                p = tf.add_paragraph()
                p.text = point
                p.level = 0
            st.warning(f"No body placeholder found for slide '{slide['title']}'. Added bullet points as a text box.")


        # Add image if available
        if slide.get('imageUrl') and not "placehold.co" in slide['imageUrl']:
            try:
                image_data = base64.b64decode(slide['imageUrl'].split(',')[1])
                img_stream = io.BytesIO(image_data)

                # Adjusted image position to reduce overlap, especially when no placeholders are found
                # This places the image in the bottom-right corner, typically a safer spot.
                img_left = Inches(7)   # Further to the right
                img_top = Inches(5.5)  # Further down
                img_width = Inches(2.5) # Smaller width
                img_height = Inches(2) # Smaller height

                pic = slide_obj.shapes.add_picture(img_stream, img_left, img_top, width=img_width, height=img_height)
            except Exception as e:
                st.warning(f"Could not add image to PPTX for slide '{slide['title']}': {e}")

    binary_output = io.BytesIO()
    prs.save(binary_output)
    binary_output.seek(0)
    return binary_output

# --- Streamlit UI ---

st.set_page_config(page_title="AI Presentation Maker", layout="centered")

st.markdown("""
    <style>
    @import url('https://fonts.googleapis.com/css2?family=Inter:wght@300;400;500;600;700&display=swap');
    html, body, [class*="st-"] {
        font-family: 'Inter', sans-serif;
    }
    .stApp {
        background-color: #000000; /* Changed to black */
        padding: 20px;
    }
    .container {
        background-color: #ffffff;
        border-radius: 1.5rem;
        box-shadow: 0 10px 25px rgba(0, 0, 0, 0.1);
        padding: 2.5rem;
        max-width: 900px;
        margin: 20px auto;
    }
    h1 {
        text-align: center;
        font-size: 2.5rem !important; /* text-4xl */
        font-weight: 800 !important; /* font-extrabold */
        margin-bottom: 2rem !important; /* mb-8 */
        background: linear-gradient(to right, #8b5cf6, #6366f1); /* from-purple-600 to-indigo-600 */
        -webkit-background-clip: text;
        -webkit-text-fill-color: transparent;
    }
    /* FIX: Ensure text is visible on black background */
    .stTextArea label, .stSelectbox label, .stFileUploader label,
    .stRadio label, .stTextInput label, .stMarkdown p, .stMarkdown li, .stMarkdown h3, .stMarkdown h4, .stInfo, .stSuccess, .stWarning, .stError {
        color: #CCCCCC !important; /* Light gray for better contrast on black */
    }
    .stTextArea textarea, .stSelectbox select, .stFileUploader input[type="file"], .stTextInput input[type="text"] {
        padding: 0.75rem; /* p-3 */
        border: 1px solid #d1d5db; /* border-gray-300 */
        border-radius: 0.5rem; /* rounded-lg */
        outline: none;
        box-shadow: none;
        color: #333333; /* Ensure input text is dark */
    }
    .stButton > button {
        background: linear-gradient(to right, #8b5cf6, #6366f1); /* from-purple-600 to-indigo-600 */
        color: white;
        font-weight: 600; /* font-semibold */
        padding: 0.75rem 1.5rem; /* py-3 px-6 */
        border-radius: 0.75rem; /* rounded-xl */
        box-shadow: 0 4px 6px rgba(0, 0, 0, 0.1); /* shadow-lg */
        transition: all 0.3s ease;
        width: 100%;
        display: flex;
        align-items: center;
        justify-content: center;
    }
    .stButton > button:hover {
        background: linear-gradient(to right, #7c3aed, #4f46e5); /* hover:from-purple-700 hover:to-indigo-700 */
        transform: scale(1.02); /* slight scale for hover effect */
    }
    .slide-card {
        background-color: #ffffff;
        border: 1px solid #e5e7eb; /* border-gray-200 */
        border-radius: 0.75rem; /* rounded-xl */
        box-shadow: 0 4px 6px rgba(0, 0, 0, 0.1); /* shadow-lg */
        padding: 1.5rem; /* p-6 */
        margin-bottom: 1.5rem; /* mb-6 */
    }
    .slide-card h3 {
        font-size: 1.25rem; /* text-xl */
        font-weight: 700; /* font-bold */
        color: #1f2937; /* gray-800 */
        margin-bottom: 0.75rem; /* mb-3 */
    }
    .slide-card ul {
        list-style-type: disc;
        padding-left: 1.25rem; /* list-inside */
        color: #374151; /* gray-700 */
    }
    .slide-card img {
        width: 100%;
        height: 12rem; /* h-48 */
        object-fit: cover;
        border-radius: 0.5rem; /* rounded-lg */
        margin-top: 1rem; /* mt-4 */
    }
    </style>
""", unsafe_allow_html=True)

st.title("AI Presentation Maker")

st.markdown("### Input Content")
input_method = st.radio(
    "Choose your input method:",
    ("Text Input", "Upload Document", "Upload Image", "YouTube Link"), # Added YouTube Link
    key="input_method",
    horizontal=True
)

content_text = ""
uploaded_file = None
youtube_url = ""

if input_method == "Text Input":
    content_text = st.text_area(
        "Enter your presentation topic or content:",
        placeholder="e.g., 'The Future of AI in Finance' or paste your full content here...",
        height=200,
        key="presentation_topic_text"
    )
elif input_method == "Upload Document":
    uploaded_file = st.file_uploader("Upload your document (TXT, DOCX, PDF):", type=["txt", "docx", "pdf"], key="file_upload")
    if uploaded_file is not None:
        file_type = uploaded_file.type
        try:
            if "text/plain" in file_type:
                content_text = read_txt(uploaded_file)
            elif "application/vnd.openxmlformats-officedocument.wordprocessingml.document" in file_type:
                content_text = read_docx(uploaded_file)
            elif "application/pdf" in file_type:
                content_text = read_pdf(uploaded_file)
            st.success("File uploaded and content read.")
        except Exception as e:
            st.error(f"Error reading file: {e}. Please ensure it's a valid TXT, DOCX, or PDF.")
            content_text = ""
    else:
        st.info("Upload a TXT, DOCX, or PDF file to provide content.")
elif input_method == "Upload Image":
    uploaded_file = st.file_uploader("Upload an image (JPG, PNG):", type=["jpg", "jpeg", "png"], key="image_upload")
    if uploaded_file is not None:
        st.image(uploaded_file, caption="Uploaded Image", use_column_width=True)
        # Optional: Allow user to add additional text context for the image
        additional_text = st.text_input("Add any additional context for the image (optional):", key="image_context_text")
        st.info("Image uploaded. AI will analyze it for presentation content.")
    else:
        st.info("Upload an image to generate presentation content based on it.")
elif input_method == "YouTube Link": # New YouTube input method
    youtube_url = st.text_input(
        "Enter YouTube Video URL:",
        placeholder="e.g., https://www.youtube.com/watch?v=dQw4w9WgXcQ",
        key="youtube_url_input"
    )
    if youtube_url:
        st.info("YouTube URL provided. The application will attempt to fetch its transcript.")


st.markdown("### Presentation Options")
col1, col2, col3 = st.columns(3)
with col1:
    num_slides = st.selectbox(
        "Number of Slides:",
        (1,2,3,4,5, 7, 10, 12, 15,20,25,30),
        index=2, # Default to 10 slides
        help="Select the desired number of slides for your presentation. Options 12 and 15 are considered 'Premium' and 'Advanced' features."
    )
with col2:
    language = st.selectbox("Language:", ("English", "Hindi", "Marathi", "Gujrati"))
with col3:
    tone = st.selectbox("Tone:", ("General", "Professional", "Informative", "Persuasive", "Casual"))

col4, col5, col6 = st.columns(3) # Added a new column for theme
with col4:
    audience = st.selectbox("Audience:", ("General", "Experts", "Beginners", "Investors"))
with col5:
    scene = st.selectbox("Scene/Style:", ("General", "Business", "Academic", "Creative"))
with col6: # New theme selection
    theme_style = st.selectbox(
        "Presentation Theme:",
        ("Default", "Classic", "Minimalist", "Organic", "Tech Design", "Basic"), # Added new themes
        index=0, # Default to "Default"
        help="Choose a visual theme for your presentation. (Requires corresponding .pptx template files)"
    )


if st.button("Generate Presentation"):
    final_topic_for_llm = ""
    image_bytes_for_llm = None

    if input_method == "Text Input":
        final_topic_for_llm = content_text
    elif input_method == "Upload Document":
        final_topic_for_llm = content_text
    elif input_method == "Upload Image" and uploaded_file is not None:
        image_bytes_for_llm = uploaded_file.getvalue()
        # Combine image description with any additional text context
        if 'image_context_text' in st.session_state and st.session_state.image_context_text:
            final_topic_for_llm = st.session_state.image_context_text
        else:
            final_topic_for_llm = "Generate a presentation based on the uploaded image."
    elif input_method == "YouTube Link" and youtube_url: # Handle YouTube Link input
        with st.spinner("Fetching YouTube transcript..."):
            transcript_content = get_youtube_transcript(youtube_url)
            if transcript_content:
                final_topic_for_llm = transcript_content
                st.success("YouTube transcript fetched successfully!")
            else:
                st.error("Could not fetch YouTube transcript. Please ensure the video has public captions/subtitles or try another video.")
                st.stop()
    
    if not final_topic_for_llm and image_bytes_for_llm is None:
        st.error("Please provide content via text, document, image upload, or a YouTube link.")
        st.stop() # Stop execution if no content is provided

    with st.spinner("Generating presentation content and images... This may take a moment."):
        # If an image is uploaded, get its description first
        if image_bytes_for_llm:
            st.info("Analyzing image content...")
            image_description = get_image_description(image_bytes_for_llm)
            if image_description:
                if final_topic_for_llm:
                    final_topic_for_llm = f"{final_topic_for_llm}\n\nImage Content: {image_description}"
                else:
                    final_topic_for_llm = f"Presentation based on: {image_description}"
            else:
                st.warning("Could not get a description for the image. Proceeding with text context only.")

        if not final_topic_for_llm: # Fallback if image description failed and no text was provided
            st.error("No content could be derived from the input. Please provide more clear input.")
            st.stop()


        slides_data = asyncio.run(generate_slide_content(final_topic_for_llm, num_slides, language, tone, audience, scene))

        if slides_data:
            st.success("Presentation content generated! Now generating images...")
            slides_with_images = []
            for i, slide in enumerate(slides_data):
                st.info(f"Generating image for slide {i+1}: '{slide['title']}'...")
                image_url = generate_image_for_slide(slide['title'])
                slides_with_images.append({"title": slide['title'], "bullet_points": slide['bullet_points'], "imageUrl": image_url})
                import time
                time.sleep(0.5) # Small delay to manage API rate limits

            st.markdown("---")
            st.markdown("## Your Generated Presentation")

            for i, slide in enumerate(slides_with_images):
                st.markdown(f"""
                    <div class="slide-card">
                        <h3>Slide {i+1}: {slide['title']}</h3>
                        <ul>
                            {"".join([f"<li>{bp}</li>" for bp in slide['bullet_points']])}
                        </ul>
                        <img src="{slide['imageUrl']}" alt="Image for {slide['title']}">
                    </div>
                """, unsafe_allow_html=True)
            st.success("Presentation generated successfully!")

            # --- Download Button ---
            st.markdown("---")
            st.markdown("### Download Presentation")
            try:
                # Pass the selected theme_style to the PPTX creation function
                pptx_file = create_pptx_file(slides_with_images, theme_style)
                st.download_button(
                    label="Download PPTX",
                    data=pptx_file,
                    file_name="generated_presentation.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    help="Click to download the generated presentation as a PowerPoint file."
                )
            except Exception as e:
                st.error(f"Error creating PPTX file: {e}. Please try again.")
        else:
            st.error("Could not generate presentation content. Please try again with a different topic or more detailed content.")

st.markdown("""
    <p class="text-center text-gray-600 mt-8 text-sm">
        *Note: This application generates slide content and relevant images. Actual .pptx file download functionality would require a dedicated Python library like `python-pptx` and a download mechanism.
    </p>
""", unsafe_allow_html=True)